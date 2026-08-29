#!/usr/bin/env python3
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import subprocess
import tarfile
import zipfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any
import xml.etree.ElementTree as ET

import fitz
import openpyxl
import pdfplumber
from PIL import Image
from docx import Document
from pptx import Presentation

try:
    import extract_msg
except ImportError:  # pragma: no cover - environment-dependent
    extract_msg = None


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
ZIP_MEDIA_PREFIXES = {
    "pptx": "ppt/media/",
    "docx": "word/media/",
    "xlsx": "xl/media/",
}


def slug_for_path(path: Path) -> str:
    digest = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:12]
    base = re.sub(r"[^a-zA-Z0-9]+", "-", path.stem).strip("-").lower() or "document"
    return f"{base[:40]}-{digest}"


def normalize_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def run_tesseract(image_path: Path, lang: str) -> str:
    cmd = ["tesseract", str(image_path), "stdout", "-l", lang]
    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
    if result.returncode != 0:
        return ""
    return normalize_text(result.stdout)


def save_image(blob: bytes, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(blob)


def extract_zip_media(source: Path, kind: str, media_dir: Path, ocr_lang: str) -> list[dict[str, Any]]:
    extracted = []
    prefix = ZIP_MEDIA_PREFIXES.get(kind)
    if not prefix:
        return extracted
    with zipfile.ZipFile(source) as zf:
        for name in zf.namelist():
            suffix = Path(name).suffix.lower()
            if not name.startswith(prefix) or suffix not in IMAGE_EXTENSIONS:
                continue
            blob = zf.read(name)
            out_path = media_dir / Path(name).name
            save_image(blob, out_path)
            ocr_text = run_tesseract(out_path, ocr_lang) if shutil_which("tesseract") else ""
            extracted.append(
                {
                    "source_member": name,
                    "saved_path": str(out_path),
                    "ocr_text": ocr_text,
                }
            )
    return extracted


def shutil_which(command: str) -> bool:
    from shutil import which

    return which(command) is not None


def extract_pdf(source: Path, out_dir: Path, ocr_lang: str) -> dict[str, Any]:
    render_dir = out_dir / "renders"
    pages = []
    full_text_parts = []
    pdf = fitz.open(source)
    with pdfplumber.open(source) as plumber_pdf:
        for index, plumber_page in enumerate(plumber_pdf.pages, start=1):
            text = normalize_text(plumber_page.extract_text() or "")
            ocr_text = ""
            rendered_path = ""
            if len(text) < 30 and shutil_which("tesseract"):
                pix = pdf[index - 1].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                render_dir.mkdir(parents=True, exist_ok=True)
                render_path = render_dir / f"page-{index:04d}.png"
                pix.save(str(render_path))
                rendered_path = str(render_path)
                ocr_text = run_tesseract(render_path, ocr_lang)
                if ocr_text and not text:
                    text = ocr_text
            pages.append(
                {
                    "page": index,
                    "text": text,
                    "ocr_text": ocr_text,
                    "rendered_path": rendered_path,
                }
            )
            full_text_parts.append(f"[[PAGE {index}]]\n{text}\n")

    payload = {
        "kind": "pdf",
        "source_path": str(source),
        "page_count": len(pages),
        "pages": pages,
        "full_text": "\n".join(full_text_parts).strip(),
    }
    write_json(out_dir / "meta.json", payload)
    (out_dir / "full_text.txt").write_text(payload["full_text"], encoding="utf-8")
    return payload


def excel_cells(ws) -> list[dict[str, Any]]:
    cells = []
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is None and cell.comment is None:
                continue
            value = "" if cell.value is None else normalize_text(str(cell.value))
            cells.append(
                {
                    "coordinate": cell.coordinate,
                    "value": value,
                    "comment": normalize_text(cell.comment.text) if cell.comment else "",
                    "data_type": cell.data_type,
                }
            )
    return cells


def extract_xlsx(source: Path, out_dir: Path, ocr_lang: str) -> dict[str, Any]:
    wb = openpyxl.load_workbook(source, data_only=False)
    sheets = []
    for ws in wb.worksheets:
        sheet_images = []
        for image in getattr(ws, "_images", []):
            anchor = getattr(image.anchor, "_from", None)
            sheet_images.append(
                {
                    "anchor_row": getattr(anchor, "row", None),
                    "anchor_col": getattr(anchor, "col", None),
                }
            )
        sheets.append(
            {
                "title": ws.title,
                "sheet_state": ws.sheet_state,
                "max_row": ws.max_row,
                "max_col": ws.max_column,
                "merged_ranges": [str(rng) for rng in ws.merged_cells.ranges],
                "images": sheet_images,
                "cells": excel_cells(ws),
            }
        )

    media = extract_zip_media(source, "xlsx", out_dir / "media", ocr_lang)
    payload = {
        "kind": "xlsx",
        "source_path": str(source),
        "sheet_count": len(sheets),
        "sheets": sheets,
        "embedded_media": media,
    }
    write_json(out_dir / "meta.json", payload)
    lines = []
    for sheet in sheets:
        lines.append(f"# Sheet: {sheet['title']} ({sheet['sheet_state']})")
        lines.append("")
        for cell in sheet["cells"]:
            line = f"{cell['coordinate']}: {cell['value']}"
            if cell["comment"]:
                line += f" [comment: {cell['comment']}]"
            lines.append(line)
        lines.append("")
    (out_dir / "full_text.txt").write_text("\n".join(lines), encoding="utf-8")
    return payload


def parse_docx_comments(source: Path) -> list[dict[str, str]]:
    comments = []
    with zipfile.ZipFile(source) as zf:
        if "word/comments.xml" not in zf.namelist():
            return comments
        root = ET.fromstring(zf.read("word/comments.xml"))
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    for comment in root.findall("w:comment", ns):
        texts = [node.text or "" for node in comment.findall(".//w:t", ns)]
        comments.append(
            {
                "id": comment.attrib.get(f"{{{ns['w']}}}id", ""),
                "author": comment.attrib.get(f"{{{ns['w']}}}author", ""),
                "text": normalize_text("".join(texts)),
            }
        )
    return comments


def parse_docx_revisions(source: Path) -> dict[str, int]:
    with zipfile.ZipFile(source) as zf:
        root = ET.fromstring(zf.read("word/document.xml"))
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    return {
        "insertions": len(root.findall(".//w:ins", ns)),
        "deletions": len(root.findall(".//w:del", ns)),
    }


def extract_docx(source: Path, out_dir: Path, ocr_lang: str) -> dict[str, Any]:
    doc = Document(source)
    paragraphs = [normalize_text(p.text) for p in doc.paragraphs if normalize_text(p.text)]
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            row_values = [normalize_text(cell.text) for cell in row.cells]
            rows.append(row_values)
        tables.append(rows)

    headers = []
    footers = []
    for section in doc.sections:
        headers.extend(normalize_text(p.text) for p in section.header.paragraphs if normalize_text(p.text))
        footers.extend(normalize_text(p.text) for p in section.footer.paragraphs if normalize_text(p.text))

    comments = parse_docx_comments(source)
    revisions = parse_docx_revisions(source)
    media = extract_zip_media(source, "docx", out_dir / "media", ocr_lang)

    payload = {
        "kind": "docx",
        "source_path": str(source),
        "paragraphs": paragraphs,
        "tables": tables,
        "headers": headers,
        "footers": footers,
        "comments": comments,
        "revisions": revisions,
        "embedded_media": media,
    }
    write_json(out_dir / "meta.json", payload)
    full_text = "\n\n".join(
        [
            "## Paragraphs\n" + "\n".join(paragraphs),
            "## Headers\n" + "\n".join(headers),
            "## Footers\n" + "\n".join(footers),
        ]
    ).strip()
    (out_dir / "full_text.txt").write_text(full_text, encoding="utf-8")
    return payload


def slide_shape_texts(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = normalize_text(shape.text)
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = normalize_text(cell.text)
                    if text:
                        texts.append(text)
    return texts


def slide_notes_text(slide) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    texts = []
    for shape in notes_slide.shapes:
        if hasattr(shape, "text"):
            text = normalize_text(shape.text)
            if text:
                texts.append(text)
    return "\n".join(texts)


def extract_pptx(source: Path, out_dir: Path, ocr_lang: str) -> dict[str, Any]:
    prs = Presentation(source)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = slide_shape_texts(slide)
        notes = slide_notes_text(slide)
        slides.append(
            {
                "slide": index,
                "texts": texts,
                "notes": notes,
            }
        )

    media = extract_zip_media(source, "pptx", out_dir / "media", ocr_lang)
    payload = {
        "kind": "pptx",
        "source_path": str(source),
        "slide_count": len(slides),
        "slides": slides,
        "embedded_media": media,
    }
    write_json(out_dir / "meta.json", payload)
    parts = []
    for slide in slides:
        parts.append(f"[[SLIDE {slide['slide']}]]")
        if slide["texts"]:
            parts.append("\n".join(slide["texts"]))
        if slide["notes"]:
            parts.append("[NOTES]")
            parts.append(slide["notes"])
        parts.append("")
    (out_dir / "full_text.txt").write_text("\n".join(parts), encoding="utf-8")
    return payload


def extract_image(source: Path, out_dir: Path, ocr_lang: str) -> dict[str, Any]:
    image = Image.open(source)
    ocr_text = run_tesseract(source, ocr_lang) if shutil_which("tesseract") else ""
    payload = {
        "kind": "image",
        "source_path": str(source),
        "size": list(image.size),
        "ocr_text": ocr_text,
    }
    write_json(out_dir / "meta.json", payload)
    (out_dir / "full_text.txt").write_text(ocr_text, encoding="utf-8")
    return payload


def unpack_archive(source: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    name = source.name.lower()
    if name.endswith(".zip"):
        with zipfile.ZipFile(source) as zf:
            zf.extractall(out_dir)
    elif name.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2")):
        with tarfile.open(source) as tf:
            tf.extractall(out_dir)
    else:
        subprocess.run(["7z", "x", "-y", f"-o{out_dir}", str(source)], check=False, capture_output=True)
    return sorted(path for path in out_dir.rglob("*") if path.is_file())


def extract_eml(source: Path, out_dir: Path) -> tuple[dict[str, Any], Path | None]:
    msg = BytesParser(policy=policy.default).parsebytes(source.read_bytes())
    attachment_dir = out_dir / "attachments"
    attachments = []
    text_parts = []
    if msg.get_body(preferencelist=("plain", "html")):
        body = msg.get_body(preferencelist=("plain", "html"))
        try:
            text_parts.append(body.get_content())
        except Exception:
            pass

    for part in msg.iter_attachments():
        filename = part.get_filename() or "attachment.bin"
        attachment_dir.mkdir(parents=True, exist_ok=True)
        saved_path = attachment_dir / filename
        payload = part.get_payload(decode=True) or b""
        saved_path.write_bytes(payload)
        attachments.append(str(saved_path))

    payload = {
        "kind": "email",
        "source_path": str(source),
        "subject": msg.get("subject", ""),
        "from": msg.get("from", ""),
        "to": msg.get("to", ""),
        "date": msg.get("date", ""),
        "body": normalize_text("\n".join(text_parts)),
        "attachments": attachments,
    }
    write_json(out_dir / "meta.json", payload)
    (out_dir / "full_text.txt").write_text(payload["body"], encoding="utf-8")
    return payload, attachment_dir if attachments else None


def extract_msg_file(source: Path, out_dir: Path) -> tuple[dict[str, Any], Path | None]:
    if extract_msg is None:
        raise RuntimeError("extract_msg is not installed")
    message = extract_msg.Message(str(source))
    attachment_dir = out_dir / "attachments"
    attachments = []
    attachment_dir.mkdir(parents=True, exist_ok=True)
    for attachment in message.attachments:
        filename = attachment.longFilename or attachment.shortFilename or "attachment.bin"
        saved_path = attachment_dir / filename
        with saved_path.open("wb") as handle:
            handle.write(attachment.data)
        attachments.append(str(saved_path))

    payload = {
        "kind": "email",
        "source_path": str(source),
        "subject": message.subject or "",
        "from": message.sender or "",
        "to": message.to or "",
        "date": str(message.date or ""),
        "body": normalize_text(message.body or ""),
        "attachments": attachments,
    }
    write_json(out_dir / "meta.json", payload)
    (out_dir / "full_text.txt").write_text(payload["body"], encoding="utf-8")
    return payload, attachment_dir if attachments else None


def process_file(source: Path, out_root: Path, ocr_lang: str) -> tuple[dict[str, Any] | None, Path | None]:
    suffix = source.suffix.lower()
    artifact_dir = out_root / "documents" / slug_for_path(source)
    artifact_dir.mkdir(parents=True, exist_ok=True)
    if suffix == ".pdf":
        return extract_pdf(source, artifact_dir, ocr_lang), None
    if suffix == ".xlsx":
        return extract_xlsx(source, artifact_dir, ocr_lang), None
    if suffix == ".docx":
        return extract_docx(source, artifact_dir, ocr_lang), None
    if suffix == ".pptx":
        return extract_pptx(source, artifact_dir, ocr_lang), None
    if suffix in IMAGE_EXTENSIONS:
        return extract_image(source, artifact_dir, ocr_lang), None
    if suffix == ".eml":
        return extract_eml(source, artifact_dir)
    if suffix == ".msg":
        return extract_msg_file(source, artifact_dir)
    if suffix in {".zip", ".7z", ".rar", ".tar", ".gz", ".tgz"} or source.name.lower().endswith(".tar.gz"):
        unpack_dir = out_root / "containers" / slug_for_path(source)
        unpack_archive(source, unpack_dir)
        payload = {
            "kind": "archive",
            "source_path": str(source),
            "expanded_to": str(unpack_dir),
        }
        write_json(artifact_dir / "meta.json", payload)
        return payload, unpack_dir

    payload = {
        "kind": "unsupported",
        "source_path": str(source),
        "note": "No native extractor configured. Consider install approval or manual review.",
    }
    write_json(artifact_dir / "meta.json", payload)
    return payload, None


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract a mixed document bundle into structured artifacts.")
    parser.add_argument("input_path", help="File or directory to extract.")
    parser.add_argument("--out", default="./_analysis", help="Output directory.")
    parser.add_argument("--ocr-lang", default="eng+chi_sim", help="Tesseract language string.")
    args = parser.parse_args()

    input_path = Path(args.input_path).resolve()
    out_root = Path(args.out).resolve()
    out_root.mkdir(parents=True, exist_ok=True)

    queue = [input_path]
    seen: set[str] = set()
    manifest = {"root": str(input_path), "documents": []}

    while queue:
        current = queue.pop(0)
        current_key = str(current.resolve())
        if current_key in seen:
            continue
        seen.add(current_key)

        if current.is_dir():
            for child in sorted(current.iterdir()):
                queue.append(child)
            continue

        meta, child_dir = process_file(current, out_root, args.ocr_lang)
        if meta is not None:
            manifest["documents"].append(meta)
        if child_dir and child_dir.exists():
            queue.append(child_dir)

    write_json(out_root / "manifest.json", manifest)
    print(out_root / "manifest.json")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
